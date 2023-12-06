from pptx import Presentation
import os

# Set your input PowerPoint file here
input_pptx = 'path_to_your_presentation.pptx'


def extract_text_from_pptx(pptx_path):
    """
    Extracts text from a PowerPoint (PPTX) file.

    :param pptx_path: Path to the PowerPoint file
    :return: Extracted text as a string
    """
    prs = Presentation(pptx_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return '\n'.join(text)


def main():
    # Set your input file here
    output_txt = os.path.splitext(input_pptx)[0] + '.txt'

    extracted_text = extract_text_from_pptx(input_pptx)

    with open(output_txt, 'w', encoding='utf-8') as file:
        file.write(extracted_text)

    print(f"Text extracted and saved to {output_txt}")


if __name__ == "__main__":
    main()
